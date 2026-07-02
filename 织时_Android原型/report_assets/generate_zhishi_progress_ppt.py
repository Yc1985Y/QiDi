from pathlib import Path
import sys

BASE_DIR = Path(__file__).resolve().parent
VENDOR_DIR = BASE_DIR / "vendor"
if str(VENDOR_DIR) not in sys.path:
    sys.path.insert(0, str(VENDOR_DIR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

OUT_FILE = BASE_DIR / "织时_项目汇报_20260609.pptx"
ICON_PATH = BASE_DIR.parent / "app" / "src" / "main" / "res" / "drawable-nodpi" / "ic_launcher_foreground_exact.png"

COLORS = {
    "bg": "F8F5EF",
    "panel": "FFFCF6",
    "panel_soft": "F2ECE2",
    "panel_blue": "EEF7F9",
    "panel_gold": "FFF4D1",
    "panel_green": "EFF7EF",
    "ink": "4A3517",
    "body": "705B33",
    "muted": "9A875E",
    "primary": "F4C64A",
    "primary_deep": "D7A82C",
    "sky": "A8D7E2",
    "sky_deep": "74BACA",
    "green": "89B68A",
    "warning": "E7B060",
    "line": "E7D9B6",
    "white": "FFFFFF",
    "soft_red": "F6E6DE",
}

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"
TOTAL_SLIDES = 14


def rgb(hex_code: str) -> RGBColor:
    hex_code = hex_code.replace("#", "")
    return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))


def set_slide_bg(slide, color_key="bg"):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(COLORS[color_key])


def add_ambient(slide, deep=False):
    set_slide_bg(slide, "panel_gold" if deep else "bg")
    for left, top, width, height, color_key, transparency in [
        (-0.2, 5.5, 3.0, 1.8, "primary", 0.82),
        (10.3, -0.15, 2.7, 1.8, "sky", 0.84),
        (11.2, 6.0, 1.1, 0.7, "primary", 0.88),
    ]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(COLORS[color_key])
        shape.fill.transparency = transparency
        shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text="", font_size=18, color="ink",
                font_name=FONT_BODY, bold=False, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    lines = text.split("\n") if text else [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = line
        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = bold
        font.color.rgb = rgb(COLORS[color] if color in COLORS else color)
        p.space_after = Pt(0)
    return box


def add_line(slide, x1, y1, x2, y2, color="line", width=1.25):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(COLORS[color] if color in COLORS else color)
    line.line.width = Pt(width)
    return line


def add_box(slide, left, top, width, height, fill="panel", line_color="white", transparency=0.0,
            radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill] if fill in COLORS else fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(COLORS[line_color] if line_color in COLORS else line_color)
    shape.line.width = Pt(1)
    return shape


def add_kicker(slide, text, left=0.72, top=0.45, width=1.35, height=0.34):
    shape = add_box(slide, left, top, width, height, fill="primary", line_color="primary")
    add_textbox(slide, left, top + 0.03, width, height - 0.05, text, font_size=10.5,
                color="white", font_name=FONT_BODY, bold=True, align=PP_ALIGN.CENTER,
                valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    return shape


def add_header(slide, kicker, title, subtitle=""):
    add_kicker(slide, kicker)
    add_textbox(slide, 0.72, 0.9, 11.0, 0.62, title, font_size=24, color="ink",
                font_name=FONT_TITLE, bold=True)
    if subtitle:
        add_textbox(slide, 0.72, 1.56, 11.2, 0.36, subtitle, font_size=11.5, color="body",
                    font_name=FONT_BODY)
    add_line(slide, 0.72, 1.98, 12.44, 1.98, color="line", width=1.3)


def add_footer(slide, page_no, total=TOTAL_SLIDES,
               source="资料依据：README / 作品描述 / 进度说明 / 配置说明"):
    add_textbox(slide, 0.72, 7.02, 8.8, 0.16, source, font_size=8.5, color="muted")
    add_textbox(slide, 11.9, 6.98, 0.7, 0.18, f"{page_no}/{total}", font_size=9,
                color="muted", font_name=FONT_MONO, bold=True, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body, fill="panel", accent="primary",
             title_size=15, body_size=11.1, mono=False):
    add_box(slide, left, top, width, height, fill=fill, line_color="white")
    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left + 0.02), Inches(top + 0.02),
        Inches(0.07), Inches(height - 0.04)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = rgb(COLORS[accent])
    accent_bar.line.fill.background()
    add_textbox(slide, left + 0.17, top + 0.14, width - 0.28, 0.3, title, font_size=title_size,
                color="ink", font_name=FONT_TITLE, bold=True)
    add_textbox(slide, left + 0.17, top + 0.48, width - 0.28, height - 0.58, body,
                font_size=body_size, color="body", font_name=FONT_MONO if mono else FONT_BODY)


def add_number_badge(slide, label, left, top, fill="primary_deep"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(0.46), Inches(0.46)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.line.fill.background()
    add_textbox(slide, left, top + 0.03, 0.46, 0.36, str(label), font_size=11,
                color="white", font_name=FONT_MONO, bold=True, align=PP_ALIGN.CENTER,
                valign=MSO_VERTICAL_ANCHOR.MIDDLE)


def add_arrow_text(slide, left, top, text="→", color="primary_deep", font_size=24):
    add_textbox(slide, left, top, 0.35, 0.35, text, font_size=font_size, color=color,
                font_name=FONT_BODY, bold=True, align=PP_ALIGN.CENTER,
                valign=MSO_VERTICAL_ANCHOR.MIDDLE)


def add_timeline_node(slide, x, y, date, title, body, up=True):
    ring = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.22), Inches(0.22)
    )
    ring.fill.solid()
    ring.fill.fore_color.rgb = rgb(COLORS["white"])
    ring.line.color.rgb = rgb(COLORS["primary"])
    ring.line.width = Pt(1.2)
    core = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.05), Inches(y + 0.05), Inches(0.12), Inches(0.12)
    )
    core.fill.solid()
    core.fill.fore_color.rgb = rgb(COLORS["primary"])
    core.line.fill.background()
    add_line(slide, x + 0.11, y if up else y + 0.22, x + 0.11, y - 0.42 if up else y + 0.64, color="line", width=1)
    top = y - 1.45 if up else y + 0.32
    accent = "primary" if up else "sky_deep"
    add_card(slide, x - 0.64, top, 1.58, 1.12, date, f"{title}\n{body}", fill="panel", accent=accent,
             title_size=12.2, body_size=9.4)


def slide01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide, deep=True)
    add_box(slide, 0.72, 0.82, 5.65, 5.82, fill="panel", line_color="white")
    add_textbox(slide, 1.05, 1.18, 2.3, 0.52, "织时", font_size=30, color="ink",
                font_name=FONT_TITLE, bold=True)
    add_textbox(slide, 1.05, 1.78, 2.3, 0.24, "TimeWeaver", font_size=12.5,
                color="primary_deep", font_name=FONT_MONO, bold=True)
    add_textbox(slide, 1.05, 2.36, 4.2, 0.42, "项目阶段汇报", font_size=22,
                color="ink", font_name=FONT_TITLE, bold=True)
    add_textbox(slide, 1.05, 2.92, 4.6, 0.26, "从设计思维流程到项目推进全过程", font_size=12.5,
                color="body", font_name=FONT_BODY)
    add_textbox(slide, 1.05, 3.56, 4.5, 0.26, "将校园信息碎片整理成清晰时间线", font_size=15.5,
                color="primary_deep", font_name=FONT_BODY, bold=True)
    add_card(
        slide, 1.02, 4.14, 4.78, 1.28, "本次汇报关注三件事",
        "• 我们为什么要做“信息降噪”\n• 设计如何转成可落地的交互结构\n• 工程如何把大模型能力约束为可信闭环",
        fill="panel_gold", accent="primary", title_size=14.2, body_size=10.5
    )
    add_textbox(slide, 1.06, 6.18, 3.1, 0.18, "团队：________________", font_size=10.8, color="muted")
    slide.shapes.add_picture(str(ICON_PATH), Inches(7.65), Inches(1.2), width=Inches(4.2), height=Inches(4.2))
    add_card(
        slide, 7.22, 5.56, 5.1, 0.92, "汇报定位",
        "这是一份“过程型项目汇报”，重点展示问题洞察、产品定义、架构设计、迭代修复与当前交付状态。",
        fill="panel_blue", accent="sky_deep", title_size=13.3, body_size=10.2
    )
    add_footer(slide, 1)


def slide02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "AGENDA", "本次汇报围绕“问题—方案—落地—进展”四条线展开",
               "不再把《织时》当作单点功能，而是把它作为一个逐步收口成型的产品系统来讲述。")
    items = [
        ("01", "设计思维起点", "从校园信息焦虑、跨 App 录入成本和认知负担出发，明确真正要解决的问题。", "panel_gold"),
        ("02", "产品定义收敛", "把“能识别”推进为“能理解、能确认、能执行、能沉淀”的 Visual-to-Tool 智能体。", "panel"),
        ("03", "设计原则与界面", "Today-First、信息降噪、Bento UI、人在回路，形成三页主结构和交互边界。", "panel_blue"),
        ("04", "工程架构与风控", "多模态输入、OCR+LLM、结构化 JSON、RiskPolicyEngine、系统 Intent 调度。", "panel"),
        ("05", "项目推进与修复", "从 API 接入、账号闭环到多事项识别、提醒语义和“我的页”分层持续迭代。", "panel_gold"),
        ("06", "当前成果与后续", "总结已经打通的主链路、仍需继续优化的边界，以及后续演示/实机推进方向。", "panel_green"),
    ]
    idx = 0
    for r in range(2):
        for c in range(3):
            n, t, b, fill = items[idx]
            x = 0.88 + c * 4.1
            y = 2.36 + r * 2.0
            accent = "sky_deep" if c == 1 else "primary"
            badge = "sky_deep" if c == 1 else "primary_deep"
            add_card(slide, x, y, 3.58, 1.55, t, b, fill=fill, accent=accent, title_size=15, body_size=10.4)
            add_number_badge(slide, n, x + 2.94, y + 0.14, fill=badge)
            idx += 1
    add_footer(slide, 2)


def slide03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "DESIGN THINKING", "设计思维让项目从“做功能”转向“解决真实场景里的混乱输入”",
               "《织时》的形成不是先画页面，而是先把校园通知场景里的用户摩擦梳理清楚。")
    add_line(slide, 1.05, 4.06, 12.15, 4.06, color="line", width=1.4)
    steps = [
        ("共情\nEmpathize", "观察群通知、讲座海报、截图和口头提醒。\n发现用户最痛苦的是“知道有事，但没法立刻转成安排”。", "panel_gold", "primary"),
        ("定义\nDefine", "问题不只是“识别文本”。\n真正的断层发生在：碎片信息无法稳定变成时间动作。", "panel", "sky_deep"),
        ("构思\nIdeate", "提出 Visual-to-Tool。\n让 AI 适应无序输入，而不是继续要求用户学表单规则。", "panel_blue", "sky_deep"),
        ("原型\nPrototype", "用 Today-First、Bento UI、确认卡和三页主壳层，把理念落成一个可操作 App。", "panel", "primary"),
        ("测试\nTest", "围绕 API 接入、交互闭环、提醒语义、异常边界与实机表现持续修复和验证。", "panel_green", "green"),
    ]
    for i, (t, b, fill, accent) in enumerate(steps):
        x = 0.82 + i * 2.48
        add_card(slide, x, 2.55, 2.02, 2.38, t, b, fill=fill, accent=accent, title_size=13.4, body_size=9.4)
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.82), Inches(3.95), Inches(0.34), Inches(0.34))
        ring.fill.solid()
        ring.fill.fore_color.rgb = rgb(COLORS["white"])
        ring.line.color.rgb = rgb(COLORS[accent])
        core = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.91), Inches(4.04), Inches(0.16), Inches(0.16))
        core.fill.solid()
        core.fill.fore_color.rgb = rgb(COLORS[accent])
        core.line.fill.background()
    add_footer(slide, 3)


def slide04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "PROBLEM", "校园通知的问题不是“看不懂”，而是“看懂以后仍然要重复劳动”",
               "传统流程把用户卡在跨应用切换和手动再录入里，《织时》要做的是把这段操作压缩到最短。")
    add_card(slide, 0.92, 2.32, 5.72, 3.8, "传统路径：信息知道了，但秩序没有建立",
             "1. 看懂海报 / 通知 / 截图\n2. 记住时间、地点、要求\n3. 再打开日历、地图或备忘录\n4. 手动录入并重新检查\n5. 忘记、漏办、错记仍然高频发生",
             fill="soft_red", accent="warning", title_size=15.2, body_size=13)
    add_card(slide, 6.84, 2.32, 5.56, 3.8, "织时路径：把无序输入压缩成一次可确认执行",
             "1. 拍照 / 分享 / 粘贴 / 语音输入\n2. AI 提取结构化事件和建议动作\n3. 用户在确认卡上完成最后授权\n4. 系统写入日历、发起导航、安排提醒\n5. 结果沉淀到个人时间线与本地资产",
             fill="panel_blue", accent="sky_deep", title_size=15.2, body_size=13)
    add_arrow_text(slide, 6.2, 3.78, "→", color="primary_deep", font_size=26)
    add_textbox(slide, 1.22, 6.36, 10.4, 0.24,
                "AI 去适应用户的无序输入，而不是让用户先学系统规则。", font_size=13.2,
                color="primary_deep", font_name=FONT_BODY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)


def slide05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "DEFINITION", "《织时》把 AI 从聊天工具改造成了一个 Visual-to-Tool 智能体中间层",
               "它不是单点 OCR，也不是自由对话，而是把输入、理解、风控和系统执行串成完整闭环。")
    add_card(slide, 0.9, 2.28, 3.9, 1.8, "多模态收敛",
             "拍照、相册、系统分享、粘贴文本、语音输入统一进入同一条理解链路。",
             fill="panel_gold", accent="primary")
    add_card(slide, 4.96, 2.28, 3.9, 1.8, "结构化意图输出",
             "大模型被约束为 strict JSON：action + confidence + payload，而不是发散式自然语言。",
             fill="panel_blue", accent="sky_deep")
    add_card(slide, 9.02, 2.28, 3.4, 1.8, "执行前风险闸门",
             "RiskPolicyEngine + HITL 负责把“能识别”提升为“敢执行”。",
             fill="panel_green", accent="green")
    add_card(slide, 0.92, 4.38, 5.24, 1.62, "一句话定位",
             "面向高校场景的智能体原型：把校园通知、海报、截图和语音，重构为可确认、可提醒、可回看的个人时间线。",
             fill="panel", accent="primary", body_size=11.3)
    add_card(slide, 6.42, 4.38, 5.98, 1.62, "不是这些，而是这些",
             "不是：只会识别图片的 Demo / 传统手填日历 / 无边界聊天助手\n而是：校园信息降噪器 / 时间秩序重构器 / 系统能力调度中间层",
             fill="panel", accent="sky_deep", body_size=11.1)
    add_footer(slide, 5)


def slide06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "PRINCIPLES", "Today-First、信息降噪与先解释再执行共同塑造了产品体验基线",
               "设计不是装饰层，而是把复杂 AI 链路压缩成用户一眼就能做决定的界面秩序。")
    add_card(slide, 0.9, 2.35, 4.0, 1.45, "Today-First",
             "首页优先展示“下一件最重要的事”，而不是先让用户面对复杂日历。", fill="panel_gold", accent="primary")
    add_card(slide, 5.05, 2.35, 3.55, 1.45, "信息降噪",
             "只保留时间、地点、事件和建议动作，把解释性噪音从首屏剥离。", fill="panel", accent="sky_deep")
    add_card(slide, 8.78, 2.35, 3.62, 1.45, "卡片化表达",
             "Bento UI 为不同任务建立清晰边界，让输入、确认、沉淀各自归位。", fill="panel_blue", accent="sky_deep")
    add_card(slide, 0.9, 4.02, 5.32, 1.48, "先解释，再执行",
             "AI 不越权代办，中高风险动作必须先展示依据，再由用户确认。", fill="panel", accent="primary")
    add_card(slide, 6.46, 4.02, 5.94, 1.48, "正常 App 化",
             "保持首页 / 时间线 / 我的三页结构，避免产品停留在工具演示感。", fill="panel_green", accent="green")
    add_textbox(slide, 0.96, 6.05, 11.0, 0.22,
                "视觉风格基线：Material Design 3 × Bento UI × Glassmorphism × 低饱和暖色系",
                font_size=11.1, color="primary_deep", font_name=FONT_BODY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)


def slide07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "FLOW", "完整主链路已经从多源导入打通到系统执行与时间线沉淀",
               "产品价值不在某一个识别动作，而在“输入 → 理解 → 决策 → 执行 → 沉淀”的整条链路。")
    steps = [
        ("01  多源收敛", "拍照 / 相册 / 分享 / 粘贴 / 语音"),
        ("02  统一建模", "CampusNoticeInput 屏蔽底层来源差异"),
        ("03  OCR + LLM", "图片先 OCR，文本与 OCR 结果再送结构化理解"),
        ("04  结构化输出", "action + confidence + payload + fallback_query"),
        ("05  风险判定", "字段完整性、阈值与动作等级进入 Policy Engine"),
        ("06  确认执行", "首页确认卡 + 日历 / 地图 / 提醒 / TTS 调度"),
        ("07  本地沉淀", "时间线、提醒偏好、导出与资产回看"),
    ]
    for i, (title, body) in enumerate(steps):
        row = 0 if i < 4 else 1
        col = i if row == 0 else i - 4
        x = 0.92 + col * 3.08 if row == 0 else 2.48 + col * 3.08
        y = 2.42 if row == 0 else 4.55
        fill = "panel_gold" if i % 2 == 0 else "panel_blue"
        accent = "primary" if i % 2 == 0 else "sky_deep"
        add_card(slide, x, y, 2.55, 1.46, title, body, fill=fill, accent=accent, title_size=12.6, body_size=10)
        if row == 0 and i < 3:
            add_arrow_text(slide, x + 2.6, 2.83, "→", color="primary_deep", font_size=24)
        if row == 1 and 4 <= i < 6:
            add_arrow_text(slide, x + 2.6, 4.96, "→", color="sky_deep", font_size=24)
    add_line(slide, 10.25, 3.9, 3.3, 4.45, color="line", width=1.1)
    add_footer(slide, 7)


def slide08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "ARCHITECTURE", "系统架构把大模型能力约束进可维护、可验证的工程边界",
               "这也是《织时》能够从概念原型走向成熟参赛作品的关键原因。")
    layers = [
        ("展示层", "首页 / 时间线 / 我的 / 启动页"),
        ("提醒与沉淀层", "ReminderScheduler / DataStore / 导出能力"),
        ("执行层", "IntentDispatcher / Calendar / Map / TTS"),
        ("风险决策层", "RiskPolicyEngine / ActionValidator / fusedConfidence"),
        ("智能理解层", "VLMNetworkClient / OCR / ResponseInterpreter"),
        ("预处理层", "图片读取 / Base64 / 文本归一化"),
        ("输入层", "相机 / 相册 / 分享 / 粘贴 / 语音"),
    ]
    for i, (title, body) in enumerate(layers):
        fill = "panel_gold" if i % 2 == 0 else "panel_blue"
        accent = "primary" if i % 2 == 0 else "sky_deep"
        add_card(slide, 0.96 + i * 0.12, 2.18 + i * 0.56, 5.85 - i * 0.24, 0.8, title, body,
                 fill=fill, accent=accent, title_size=12.2, body_size=10)
    add_card(slide, 7.15, 2.3, 5.05, 1.02, "关键代码模块",
             "MainActivity · HomeScreenModule · TimelineScreenModule · ProfileScreenModule",
             fill="panel", accent="primary", body_size=10.2, mono=True)
    add_card(slide, 7.15, 3.52, 5.05, 1.02, "模型与接口",
             "VLM_APP_ID / VLM_API_KEY 已接入 vivo 真接口\n文本直连 chat/completions；图片走 OCR + LLM 双阶段链路",
             fill="panel_blue", accent="sky_deep", body_size=10.1)
    add_card(slide, 7.15, 4.74, 5.05, 1.02, "工程边界",
             "不回退 mock、不破坏页面基线、不把风控移除、不把执行链重新做成黑盒。",
             fill="panel_green", accent="green", body_size=10.1)
    add_footer(slide, 8)


def slide09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "SAFETY", "风险决策与人在回路是项目从 Demo 走向产品的关键分水岭",
               "《织时》不是把模型结果直接映射成动作，而是让客户端保留最后一层解释、追问和拦截。")
    add_card(slide, 0.94, 2.32, 5.1, 3.78, "Extract → Suggest → Confirm → Execute",
             "Extract：获取结构化 JSON 结果\n\nSuggest：校验字段完整性、置信度和动作等级\n\nConfirm：中高风险动作必须进入确认或澄清\n\nExecute：只把被允许的动作映射给系统 Intent",
             fill="panel", accent="primary", title_size=15, body_size=11.3)
    add_card(slide, 6.32, 2.32, 6.0, 1.22, "当前动作风险分级",
             "send_sms：高风险，绝不自动发送    create_event / navigate：中风险，必须确认    tts_feedback：低风险，可直接播报",
             fill="panel_gold", accent="warning", title_size=14.1, body_size=10.5)
    add_card(slide, 6.32, 3.76, 6.0, 1.02, "结构化契约",
             "action + confidence + payload + fallback_query + target_found",
             fill="panel_blue", accent="sky_deep", body_size=11, mono=True)
    add_card(slide, 6.32, 5.0, 6.0, 1.1, "安全价值",
             "让评委和用户都能清楚看到：系统为什么建议、为什么拦截、为什么需要再确认。",
             fill="panel_green", accent="green", body_size=11)
    add_footer(slide, 9)


def slide10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "UI / IA", "三页主结构让高频操作、资产沉淀和控制中枢各归其位",
               "视觉上保持 Minimalist Tech + Bento UI，信息上坚持先焦点、再结构、再控制。")
    add_card(slide, 0.92, 2.28, 3.83, 3.9, "首页 Home",
             "角色：动态 AI 工作台\n\n• Today-First 焦点卡\n• 多模态意图输入舱\n• 白盒化确认卡\n• 最近识别结果与继续动作\n\n首页解决的是“现在要处理什么”。",
             fill="panel_gold", accent="primary", title_size=16, body_size=10.9)
    add_card(slide, 4.95, 2.28, 3.83, 3.9, "时间线 Timeline",
             "角色：秩序化资产沉淀区\n\n• 日 / 周 / 月聚合\n• 垂直时间轴列表\n• 提醒状态与导出\n• 日程详情与回看\n\n时间线解决的是“这些事如何被组织起来”。",
             fill="panel_blue", accent="sky_deep", title_size=16, body_size=10.9)
    add_card(slide, 8.98, 2.28, 3.42, 3.9, "我的 Profile",
             "角色：智能体控制中枢\n\n• 账号资料与偏好\n• 数字福祉与统计\n• 历史 / 账户 / 设置 / 更多功能\n• 风险开关与提醒策略\n\n这里解决的是“我如何掌控这个智能体”。",
             fill="panel", accent="green", title_size=16, body_size=10.7)
    add_footer(slide, 10)


def slide11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "MILESTONES", "项目推进遵循“先打主链路，再补系统能力，再修关键问题”的节奏",
               "时间线上的每一次收口，都是为了让作品更接近真正可演示、可答辩、可交付的状态。")
    timeline = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.1), Inches(4.09), Inches(10.9), Inches(0.03))
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = rgb(COLORS["line"])
    timeline.line.fill.background()
    nodes = [
        (1.2, 4.0, "05.11", "API / UI 基线", "真实模型接入\n三页结构定型", True),
        (3.1, 4.0, "05.12", "测试闭环", "自动化验证\n交互语义修正", False),
        (5.0, 4.0, "05.29", "语音与交付", "构建、安装、APK 输出", True),
        (6.9, 4.0, "06.01", "编译恢复", "修复编码问题\n恢复稳定构建", False),
        (8.8, 4.0, "06.06", "识别兜底", "多事项切分\n提高通知解析覆盖", True),
        (10.7, 4.0, "06.07", "产品收口", "提醒语义统一\n我的页结构补齐", False),
    ]
    for n in nodes:
        add_timeline_node(slide, *n)
    add_footer(slide, 11)


def slide12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "ISSUES & FIXES", "近期修复重点已经从“能不能跑”转向“跑得是否可信、是否顺手”",
               "这部分最能体现项目从展示型原型往成熟型参赛作品靠近的速度。")
    add_card(slide, 0.92, 2.34, 5.58, 3.78, "已经明确修掉或显著收口的问题",
             "• 分享文本和图片已真实跑通到 OCR / 模型 / 确认卡阶段\n• 多事项通知识别增加本地切分兜底，不再只拿到第一条\n• 提醒数统一按“待提醒事项数”统计，避免语义误导\n• 首页主按钮逻辑、重复登录、分享后不回首页等问题已收口\n• “我的页”从堆叠入口改为统计 / 账户 / 设置等二级结构",
             fill="panel_green", accent="green", title_size=14.8, body_size=11)
    add_card(slide, 6.72, 2.34, 5.66, 3.78, "仍需继续打磨的边界",
             "• 语音输入全链路与实机表现仍需进一步验证\n• 地图联动的真实设备适配仍需持续观察\n• 图片低置信度分支与澄清体验还可继续优化\n• 导出 JPG / PNG 与更多边界场景还需要回归测试\n• 页面文案、轻量动效和性能表现仍可继续统一",
             fill="panel", accent="warning", title_size=14.8, body_size=11)
    add_footer(slide, 12)


def slide13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "CURRENT STATUS", "当前版本已经具备可答辩、可演示、可继续扩展的完整闭环",
               "它已经不是“概念图 + 说明文档”，而是一个有真实模型、有真实系统能力、有真实风控的运行原型。")
    add_card(slide, 0.92, 2.28, 6.0, 3.9, "当前已具备的核心能力",
             "√ 多入口导入：拍照 / 相册 / 分享 / 粘贴 / 语音\n√ vivo AI 接入：文本直连 LLM，图片走 OCR + LLM\n√ 结构化 JSON 输出与解释式确认卡\n√ 系统执行：日历、地图、提醒、TTS\n√ 时间线沉淀、本地持久化与偏好控制\n√ Profile 二级入口、统计面板、账号结构\n√ 构建、安装、模拟器验证与 APK 交付路径",
             fill="panel_gold", accent="primary", title_size=14.8, body_size=11)
    add_card(slide, 7.16, 2.28, 5.22, 1.22, "环境与协作基线",
             "E 盘 Android Studio + 指定 JDK / SDK / AVD；进度说明与配置说明作为统一接手入口。",
             fill="panel_blue", accent="sky_deep", body_size=10.8)
    add_card(slide, 7.16, 3.74, 5.22, 1.22, "当前演示价值",
             "足以支撑“信息降噪 → 结构化理解 → 安全执行 → 时间线沉淀”的完整答辩叙事。",
             fill="panel_green", accent="green", body_size=10.8)
    add_card(slide, 7.16, 5.2, 5.22, 0.98, "一句话结论",
             "《织时》已经具备成熟参赛作品的产品边界和工程可信度。",
             fill="panel", accent="primary_deep", body_size=10.8)
    add_footer(slide, 13)


def slide14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_ambient(slide)
    add_header(slide, "NEXT", "下一阶段已经很清晰：围绕实机体验、语音链路与性能继续收口",
               "这页保留为后续展示与答辩过渡页，你也可以在这里补充演示视频、二维码或团队信息。")
    add_card(slide, 0.94, 2.35, 3.74, 2.2, "P0 真实闭环",
             "继续把语音、拍照、导航和提醒在真实设备上的闭环验证做扎实，确保演示时不掉链子。",
             fill="panel_gold", accent="primary", title_size=15, body_size=11.1)
    add_card(slide, 4.86, 2.35, 3.74, 2.2, "P1 体验收口",
             "继续优化低置信度澄清、页面文案一致性、交互反馈和各页流畅度，减少答辩时的感知阻力。",
             fill="panel_blue", accent="sky_deep", title_size=15, body_size=11.1)
    add_card(slide, 8.78, 2.35, 3.62, 2.2, "P2 交付表达",
             "结合实机演示视频、答辩口播和阶段汇报材料，把产品逻辑、设计语言和工程可信度讲得更完整。",
             fill="panel_green", accent="green", title_size=15, body_size=11.1)
    placeholder = add_box(slide, 0.96, 5.02, 4.9, 1.08, fill="panel", line_color="primary")
    placeholder.line.dash_style = 1
    add_textbox(slide, 1.16, 5.28, 4.45, 0.5, "演示视频 / 二维码占位\n后续可直接替换为录屏、实机视频或展示链接",
                font_size=12, color="body", font_name=FONT_BODY, align=PP_ALIGN.CENTER,
                valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_textbox(slide, 7.9, 6.36, 2.9, 0.18, "团队：________________", font_size=11.2, color="muted")
    add_textbox(slide, 6.42, 5.18, 5.8, 0.5,
                "结论：我们已经把《织时》从“想法”推进到了“可被验证的产品形态”。",
                font_size=18.2, color="ink", font_name=FONT_TITLE, bold=True)
    add_textbox(slide, 6.42, 5.82, 5.8, 0.32,
                "后续只需要继续围绕真实设备体验和边界稳定性打磨，就能让它更接近完整交付。",
                font_size=11.5, color="body")
    add_footer(slide, 14, TOTAL_SLIDES, "资料依据：配置说明.md / 进度说明.md / 作品描述.md / README.md / demo_script.md")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.author = "Codex"
    prs.core_properties.title = "《织时》项目汇报"
    prs.core_properties.subject = "从设计思维流程到项目推进全过程"
    prs.core_properties.comments = "Generated from current project documents."

    slide01(prs)
    slide02(prs)
    slide03(prs)
    slide04(prs)
    slide05(prs)
    slide06(prs)
    slide07(prs)
    slide08(prs)
    slide09(prs)
    slide10(prs)
    slide11(prs)
    slide12(prs)
    slide13(prs)
    slide14(prs)

    prs.save(OUT_FILE)
    print(f"PPTX written to: {OUT_FILE}")


if __name__ == "__main__":
    build()
